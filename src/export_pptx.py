"""PowerPoint export: generate a PPTX report with key findings and charts."""

import io
from datetime import date

import polars as pl
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.chart import XL_CHART_TYPE


# Brand colors
BLUE = RGBColor(0x4E, 0x79, 0xA7)
RED = RGBColor(0xE1, 0x57, 0x59)
GREEN = RGBColor(0x59, 0xA1, 0x4F)
ORANGE = RGBColor(0xF2, 0x8E, 0x2B)
DARK = RGBColor(0x33, 0x33, 0x33)
GRAY = RGBColor(0x88, 0x88, 0x88)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)


def _add_title_slide(prs: Presentation, title: str, subtitle: str):
    """Add a title slide."""
    layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle


def _add_section_slide(prs: Presentation, title: str):
    """Add a section divider slide."""
    layout = prs.slide_layouts[2]  # Section header
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title


def _add_content_slide(prs: Presentation, title: str, bullets: list[str]):
    """Add a slide with title and bullet points."""
    layout = prs.slide_layouts[1]  # Title + Content
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.clear()
    for i, bullet in enumerate(bullets):
        if i == 0:
            tf.text = bullet
        else:
            p = tf.add_paragraph()
            p.text = bullet
            p.level = 0


def _add_kpi_slide(prs: Presentation, title: str, kpis: list[tuple[str, str]]):
    """Add a slide with large KPI numbers."""
    layout = prs.slide_layouts[5]  # Blank
    slide = prs.slides.add_slide(layout)

    # Title
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = txBox.text_frame
    tf.text = title
    tf.paragraphs[0].font.size = Pt(28)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = DARK

    # KPI boxes
    n_kpis = len(kpis)
    box_width = min(2.2, 9.0 / max(n_kpis, 1))
    start_x = (10 - box_width * n_kpis) / 2

    for i, (label, value) in enumerate(kpis):
        x = Inches(start_x + i * box_width)
        y = Inches(1.8)
        w = Inches(box_width - 0.2)
        h = Inches(2.5)

        shape = slide.shapes.add_textbox(x, y, w, h)
        tf = shape.text_frame
        tf.word_wrap = True

        # Value
        p = tf.paragraphs[0]
        p.text = str(value)
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = BLUE
        p.alignment = PP_ALIGN.CENTER

        # Label
        p2 = tf.add_paragraph()
        p2.text = label
        p2.font.size = Pt(12)
        p2.font.color.rgb = GRAY
        p2.alignment = PP_ALIGN.CENTER


def _add_table_slide(prs: Presentation, title: str, headers: list[str], rows: list[list[str]], max_rows: int = 12):
    """Add a slide with a data table."""
    layout = prs.slide_layouts[5]  # Blank
    slide = prs.slides.add_slide(layout)

    # Title
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    tf = txBox.text_frame
    tf.text = title
    tf.paragraphs[0].font.size = Pt(24)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = DARK

    # Table
    display_rows = rows[:max_rows]
    n_rows = len(display_rows) + 1  # +1 for header
    n_cols = len(headers)

    table_shape = slide.shapes.add_table(
        n_rows, n_cols,
        Inches(0.5), Inches(1.2),
        Inches(9), Inches(min(5.0, n_rows * 0.4)),
    )
    table = table_shape.table

    # Header row
    for j, header in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = BLUE
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.color.rgb = WHITE
            paragraph.font.size = Pt(10)
            paragraph.font.bold = True

    # Data rows
    for i, row in enumerate(display_rows):
        for j, val in enumerate(row):
            cell = table.cell(i + 1, j)
            cell.text = str(val)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(9)
                paragraph.font.color.rgb = DARK
            # Alternating row color
            if i % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xF0, 0xF0, 0xF0)

    if len(rows) > max_rows:
        txBox2 = slide.shapes.add_textbox(
            Inches(0.5), Inches(6.5), Inches(9), Inches(0.3),
        )
        tf2 = txBox2.text_frame
        tf2.text = f"Showing {max_rows} of {len(rows)} rows"
        tf2.paragraphs[0].font.size = Pt(8)
        tf2.paragraphs[0].font.color.rgb = GRAY


def generate_pptx(
    message_fact: pl.DataFrame,
    edge_fact: pl.DataFrame,
    person_dim: pl.DataFrame,
    graph_metrics: pl.DataFrame,
    health_score: dict | None = None,
    bridges: pl.DataFrame | None = None,
    start_date: date | None = None,
    end_date: date | None = None,
    org_name: str = "Organization",
) -> bytes:
    """Generate a PowerPoint report.

    Returns the PPTX file as bytes.
    """
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    date_range = ""
    if start_date and end_date:
        date_range = f"{start_date.strftime('%b %d, %Y')} — {end_date.strftime('%b %d, %Y')}"

    # --- Title Slide ---
    _add_title_slide(
        prs,
        f"Email Communication Analysis",
        f"{org_name}\n{date_range}\n\nGenerated by Email Metadata Analytics Platform",
    )

    # --- Executive Summary KPIs ---
    total_msgs = len(message_fact)
    unique_people = len(person_dim)
    unique_senders = message_fact["from_email"].n_unique() if total_msgs > 0 else 0
    n_communities = graph_metrics["community_id"].n_unique() if len(graph_metrics) > 0 else 0

    _add_kpi_slide(prs, "Executive Summary", [
        ("Total Messages", f"{total_msgs:,}"),
        ("Unique People", f"{unique_people:,}"),
        ("Active Senders", f"{unique_senders:,}"),
        ("Communities", f"{n_communities:,}"),
    ])

    # --- Health Score ---
    if health_score:
        composite = health_score["composite"]
        sub_scores = health_score["sub_scores"]
        verdict = "Healthy" if composite >= 70 else ("Moderate" if composite >= 50 else "Needs Attention")
        health_bullets = [f"Overall Health Score: {composite:.0f}/100 ({verdict})"]
        for key, score in sub_scores.items():
            health_bullets.append(f"{score['label']}: {score['value']:.0f}/100 — {score['detail']}")
        _add_content_slide(prs, "Organizational Health Score", health_bullets)

    # --- Communication Volume ---
    _add_section_slide(prs, "Communication Patterns")

    if total_msgs > 0:
        # After-hours stats
        ah_rate = float(message_fact["is_after_hours"].mean()) * 100
        weekend_rate = float(message_fact["is_weekend"].mean()) * 100
        avg_size_kb = float(message_fact["size_bytes"].mean()) / 1024

        _add_kpi_slide(prs, "Work Pattern Metrics", [
            ("After-Hours Rate", f"{ah_rate:.1f}%"),
            ("Weekend Rate", f"{weekend_rate:.1f}%"),
            ("Avg Message Size", f"{avg_size_kb:.1f} KB"),
        ])

    # --- Top Senders ---
    if total_msgs > 0:
        top_senders = (
            message_fact.group_by("from_email")
            .agg(pl.len().alias("count"))
            .sort("count", descending=True)
            .head(15)
        )
        headers = ["Email", "Messages Sent"]
        rows = [
            [row["from_email"], f"{row['count']:,}"]
            for row in top_senders.iter_rows(named=True)
        ]
        _add_table_slide(prs, "Top Senders by Volume", headers, rows)

    # --- Network Structure ---
    _add_section_slide(prs, "Network Structure")

    if len(graph_metrics) > 0:
        # Top connectors (betweenness)
        top_connectors = (
            graph_metrics.sort("betweenness_centrality", descending=True)
            .head(10)
        )
        headers = ["Person", "Connector Score", "Importance", "Community"]
        rows = [
            [
                row["email"],
                f"{row['betweenness_centrality']:.4f}",
                f"{row['pagerank']:.4f}",
                str(row.get("community_label", row.get("community_id", ""))),
            ]
            for row in top_connectors.iter_rows(named=True)
        ]
        _add_table_slide(prs, "Key Connectors (Highest Betweenness)", headers, rows)

        # Community summary
        comm_summary = (
            graph_metrics.group_by("community_id")
            .agg([
                pl.len().alias("members"),
                pl.col("community_label").first().alias("label"),
            ])
            .sort("members", descending=True)
            .head(15)
        )
        headers = ["Community", "Members"]
        rows = [
            [str(row.get("label", f"Group {row['community_id']}")), f"{row['members']:,}"]
            for row in comm_summary.iter_rows(named=True)
        ]
        _add_table_slide(prs, "Largest Communication Communities", headers, rows)

    # --- Bridges ---
    if bridges is not None and len(bridges) > 0:
        top_bridges = bridges.head(10)
        headers = ["Person", "Communities Bridged"]
        rows = [
            [row["email"], str(row["communities_bridged"])]
            for row in top_bridges.iter_rows(named=True)
        ]
        _add_table_slide(prs, "Bridge People (Cross-Group Connectors)", headers, rows)

    # --- Department Breakdown ---
    if "department" in person_dim.columns:
        dept_summary = (
            person_dim.group_by("department")
            .agg(pl.len().alias("people"))
            .sort("people", descending=True)
            .head(15)
        )
        if len(dept_summary) > 1:
            headers = ["Department", "People"]
            rows = [
                [row["department"], f"{row['people']:,}"]
                for row in dept_summary.iter_rows(named=True)
            ]
            _add_table_slide(prs, "Department Breakdown", headers, rows)

    # --- Closing slide ---
    _add_content_slide(prs, "Methodology & Limitations", [
        "Analysis based on email metadata only (Date, From, To, Size)",
        "Message bodies are never accessed or stored",
        "Community detection uses Louvain/Leiden algorithm on communication graph",
        "Reply time is estimated from message timing patterns (A→B then B→A within 24h)",
        "Automated senders are detected by name patterns and send/receive ratios",
        f"Date range: {date_range}" if date_range else "Full dataset analyzed",
        "Generated by Email Metadata Analytics Platform",
    ])

    # Save to bytes
    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.getvalue()
