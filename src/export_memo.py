"""Executive memo PDF: a short, opinionated consulting deliverable."""

import io
from datetime import date

import polars as pl
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Reuse brand colors from export_pptx
BLUE = RGBColor(0x4E, 0x79, 0xA7)
DARK = RGBColor(0x33, 0x33, 0x33)
GRAY = RGBColor(0x88, 0x88, 0x88)
RED = RGBColor(0xE1, 0x57, 0x59)
GREEN = RGBColor(0x59, 0xA1, 0x4F)
ORANGE = RGBColor(0xF2, 0x8E, 0x2B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)


def _text_box(slide, left, top, width, height, text, font_size=11,
              bold=False, color=DARK, align=PP_ALIGN.LEFT):
    """Add a text box to a slide."""
    shape = slide.shapes.add_textbox(Inches(left), Inches(top),
                                     Inches(width), Inches(height))
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align
    return tf


def _add_bullet_slide(prs, title, bullets, subtitle=None):
    """Add a slide with title and bullet list."""
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank

    _text_box(slide, 0.5, 0.3, 9, 0.6, title, font_size=24, bold=True, color=BLUE)

    if subtitle:
        _text_box(slide, 0.5, 0.9, 9, 0.4, subtitle, font_size=11, color=GRAY)

    y_start = 1.5 if subtitle else 1.2
    for i, bullet in enumerate(bullets):
        severity = bullet.get("severity", "info") if isinstance(bullet, dict) else "info"
        text = bullet.get("text", bullet) if isinstance(bullet, dict) else bullet
        color = RED if severity == "critical" else ORANGE if severity == "warning" else DARK
        prefix = "\u26a0 " if severity in ("critical", "warning") else "\u2022 "

        _text_box(slide, 0.7, y_start + i * 0.45, 8.5, 0.4,
                  f"{prefix}{text}", font_size=12, color=color)


def generate_executive_memo(
    message_fact: pl.DataFrame,
    edge_fact: pl.DataFrame,
    person_dim: pl.DataFrame,
    graph_metrics: pl.DataFrame,
    health_score: dict | None = None,
    narrative: str = "",
    alerts: list[dict] | None = None,
    start_date: date | None = None,
    end_date: date | None = None,
    org_name: str = "Organization",
) -> bytes:
    """Generate a short executive memo as PPTX (saves as PDF-like deck).

    This is a 5-7 slide consulting memo, not a full report.
    """
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    date_range = ""
    if start_date and end_date:
        date_range = f"{start_date.strftime('%b %d, %Y')} \u2014 {end_date.strftime('%b %d, %Y')}"

    total_msgs = len(message_fact)
    unique_people = len(person_dim)

    # -----------------------------------------------------------------------
    # Slide 1: Cover
    # -----------------------------------------------------------------------
    cover = prs.slides.add_slide(prs.slide_layouts[5])
    # Blue bar at top
    shape = cover.shapes.add_shape(
        1, Inches(0), Inches(0), Inches(10), Inches(0.15),  # 1 = rectangle
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = BLUE
    shape.line.fill.background()

    _text_box(cover, 0.8, 2.0, 8.4, 1.2,
              "Email Communication Analysis", font_size=36, bold=True, color=DARK)
    _text_box(cover, 0.8, 3.2, 8.4, 0.5,
              f"Executive Briefing \u2014 {org_name}", font_size=18, color=BLUE)
    _text_box(cover, 0.8, 4.0, 8.4, 0.4, date_range, font_size=14, color=GRAY)
    _text_box(cover, 0.8, 6.5, 8.4, 0.3,
              "CONFIDENTIAL \u2014 Prepared by Email Metadata Analytics Platform",
              font_size=10, color=GRAY)

    # -----------------------------------------------------------------------
    # Slide 2: Key Numbers
    # -----------------------------------------------------------------------
    nums = prs.slides.add_slide(prs.slide_layouts[5])
    _text_box(nums, 0.5, 0.3, 9, 0.6, "Key Numbers", font_size=24, bold=True, color=BLUE)

    unique_senders = message_fact["from_email"].n_unique() if total_msgs > 0 else 0
    n_comms = graph_metrics["community_id"].n_unique() if len(graph_metrics) > 0 else 0
    ah_rate = float(message_fact["is_after_hours"].mean()) * 100 if total_msgs > 0 else 0

    kpis = [
        ("Messages Analyzed", f"{total_msgs:,}"),
        ("Unique People", f"{unique_people:,}"),
        ("Active Senders", f"{unique_senders:,}"),
        ("Communication Groups", f"{n_comms:,}"),
        ("After-Hours Rate", f"{ah_rate:.1f}%"),
    ]

    if health_score:
        kpis.append(("Health Score", f"{health_score['composite']:.0f}/100"))

    for i, (label, value) in enumerate(kpis):
        col = i % 3
        row = i // 3
        x = 0.8 + col * 3.0
        y = 1.5 + row * 2.2

        _text_box(nums, x, y, 2.5, 0.8, value,
                  font_size=32, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
        _text_box(nums, x, y + 0.8, 2.5, 0.3, label,
                  font_size=11, color=GRAY, align=PP_ALIGN.CENTER)

    # -----------------------------------------------------------------------
    # Slide 3: Health Assessment
    # -----------------------------------------------------------------------
    if health_score:
        composite = health_score["composite"]
        verdict = "Healthy" if composite >= 70 else ("Needs Attention" if composite < 50 else "Moderate")
        verdict_color = GREEN if composite >= 70 else (RED if composite < 50 else ORANGE)

        health_slide = prs.slides.add_slide(prs.slide_layouts[5])
        _text_box(health_slide, 0.5, 0.3, 9, 0.6,
                  "Health Assessment", font_size=24, bold=True, color=BLUE)

        _text_box(health_slide, 1.0, 1.5, 3.0, 1.0,
                  f"{composite:.0f}", font_size=64, bold=True,
                  color=verdict_color, align=PP_ALIGN.CENTER)
        _text_box(health_slide, 1.0, 2.5, 3.0, 0.4,
                  verdict, font_size=18, color=verdict_color, align=PP_ALIGN.CENTER)

        y = 1.5
        for key, score in health_score["sub_scores"].items():
            val = score["value"]
            sc = GREEN if val >= 70 else (ORANGE if val >= 50 else RED)
            _text_box(health_slide, 4.5, y, 3.0, 0.3,
                      score["label"], font_size=11, color=DARK)
            _text_box(health_slide, 7.5, y, 1.5, 0.3,
                      f"{val:.0f}/100", font_size=11, bold=True, color=sc)
            y += 0.4

    # -----------------------------------------------------------------------
    # Slide 4: Key Findings (from narrative)
    # -----------------------------------------------------------------------
    if narrative:
        findings = [line.strip() for line in narrative.split("\n\n") if line.strip()]
        # Strip markdown bold
        findings = [f.replace("**", "") for f in findings]
        _add_bullet_slide(prs, "Key Findings", findings[:8])

    # -----------------------------------------------------------------------
    # Slide 5: Risks & Alerts
    # -----------------------------------------------------------------------
    if alerts:
        # Group by severity
        critical = [a for a in alerts if a.get("severity") == "critical"]
        warnings = [a for a in alerts if a.get("severity") == "warning"]

        risk_bullets = []
        for a in (critical + warnings)[:10]:
            risk_bullets.append({
                "text": f"{a.get('name', 'Alert')}: {a.get('detail', '')}",
                "severity": a.get("severity", "info"),
            })

        if risk_bullets:
            _add_bullet_slide(prs, "Risks & Alerts",
                              risk_bullets,
                              subtitle=f"{len(critical)} critical, {len(warnings)} warnings")

    # -----------------------------------------------------------------------
    # Slide 6: Top Recommendations
    # -----------------------------------------------------------------------
    recommendations = _generate_recommendations(
        message_fact, graph_metrics, health_score, alerts,
    )
    if recommendations:
        _add_bullet_slide(prs, "Recommendations", recommendations[:6])

    # -----------------------------------------------------------------------
    # Slide 7: Methodology
    # -----------------------------------------------------------------------
    _add_bullet_slide(prs, "Methodology & Scope", [
        f"Analysis period: {date_range}" if date_range else "Full dataset analyzed",
        "Data source: email metadata only (Date, From, To, Size)",
        "Message bodies are never accessed or stored",
        "Automated senders excluded from human communication metrics",
        "Community detection via Louvain/Leiden algorithm",
        "Reply time estimated from message timing patterns",
    ])

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.getvalue()


def _generate_recommendations(
    message_fact: pl.DataFrame,
    graph_metrics: pl.DataFrame,
    health_score: dict | None,
    alerts: list[dict] | None,
) -> list[str]:
    """Auto-generate recommendations based on findings."""
    recs = []

    if health_score:
        sub = health_score.get("sub_scores", {})

        if sub.get("work_life", {}).get("value", 100) < 50:
            recs.append(
                "Review after-hours communication norms. Consider establishing "
                "expected response windows to prevent burnout."
            )

        if sub.get("reciprocity", {}).get("value", 100) < 50:
            recs.append(
                "Low reciprocity suggests one-way broadcast patterns. "
                "Encourage two-way dialogue and feedback loops."
            )

        if sub.get("resilience", {}).get("value", 100) < 50:
            recs.append(
                "Network is fragile \u2014 a few key people are critical bottlenecks. "
                "Cross-train and establish backup communication paths."
            )

        if sub.get("cross_group", {}).get("value", 100) < 50:
            recs.append(
                "Communication silos detected. Consider cross-functional meetings "
                "or liaison roles between isolated groups."
            )

    if alerts:
        critical_count = sum(1 for a in alerts if a.get("severity") == "critical")
        if critical_count > 0:
            recs.append(
                f"Address {critical_count} critical alert(s) immediately \u2014 "
                "these represent significant organizational risks."
            )

    if len(graph_metrics) > 0:
        max_bc = float(graph_metrics["betweenness_centrality"].max())
        if max_bc > 0.05:
            top_person = graph_metrics.sort("betweenness_centrality", descending=True)["email"][0]
            recs.append(
                f"Reduce single-point-of-failure risk around {top_person.split('@')[0]} "
                f"(connector score: {max_bc:.3f}). Ensure knowledge transfer."
            )

    if not recs:
        recs.append("No urgent recommendations \u2014 communication patterns appear healthy.")

    return recs
